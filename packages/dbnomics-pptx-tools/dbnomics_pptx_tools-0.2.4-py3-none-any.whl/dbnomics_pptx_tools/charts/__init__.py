from typing import Iterable, cast

import daiquiri
from pptx.chart.axis import DateAxis
from pptx.chart.chart import Chart
from pptx.chart.data import CategoryChartData
from pptx.chart.datalabel import DataLabel
from pptx.chart.plot import _BasePlot
from pptx.chart.point import Point
from pptx.chart.series import _BaseCategorySeries
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.shapetree import SlideShapes
from pptx.slide import Slide

from dbnomics_pptx_tools.metadata import ChartSpec, DataLabelPosition
from dbnomics_pptx_tools.pptx_copy import copy_shape_properties
from dbnomics_pptx_tools.repo import SeriesRepo
from dbnomics_pptx_tools.xml_utils import remove_element

from .chart_data import build_category_chart_data, load_chart_df
from .data_labels import add_data_label_to_last_point_of_each_series

__all__ = ["update_chart"]

logger = daiquiri.getLogger(__name__)


def recreate_chart(chart_shape: GraphicFrame, *, chart_data: CategoryChartData, slide: Slide):
    chart = cast(Chart, chart_shape.chart)
    remove_element(chart_shape.element)
    new_chart_shape = cast(
        GraphicFrame,
        cast(SlideShapes, slide.shapes).add_chart(
            chart.chart_type, chart_shape.left, chart_shape.top, chart_shape.width, chart_shape.height, chart_data
        ),
    )
    copy_shape_properties(chart_shape, new_chart_shape)
    logger.debug("The chart was recreated")
    return cast(Chart, new_chart_shape.chart)


def remove_data_labels(chart: Chart) -> None:
    for plot_index, plot in enumerate(cast(Iterable[_BasePlot], chart.plots)):
        if plot.has_data_labels:
            logger.debug("Plot #%d has data labels, removing", plot_index)
            plot.has_data_labels = False
        for series in cast(Iterable[_BaseCategorySeries], plot.series):
            for point_index, point in enumerate(cast(Iterable[Point], series.points)):
                data_label = cast(DataLabel, point.data_label)
                if data_label._dLbl is not None:
                    logger.debug("Point #%d of series %r has a data label, removing", point_index, series.name)
                    remove_element(data_label._dLbl)


def update_chart(chart_shape: GraphicFrame, *, chart_spec: ChartSpec, repo: SeriesRepo, slide: Slide):
    chart = cast(Chart, chart_shape.chart)
    if not isinstance(chart.category_axis, DateAxis):
        raise NotImplementedError()

    df = load_chart_df(chart_spec, repo=repo)
    chart_data = build_category_chart_data(chart_spec, df=df)

    try:
        chart.replace_data(chart_data)
    except ValueError:
        chart = recreate_chart(chart_shape, chart_data=chart_data, slide=slide)
    else:
        logger.debug("Chart data was replaced")

    remove_data_labels(chart)
    if DataLabelPosition.LAST_POINT.value in chart_spec.data_labels:
        add_data_label_to_last_point_of_each_series(chart, df=df)
